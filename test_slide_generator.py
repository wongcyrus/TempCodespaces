import unittest
import pptx
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

class TestSlideGenerator(unittest.TestCase):

    def setUp(self):
        self.prs = pptx.Presentation()

    def tearDown(self):
      if os.path.exists("test_presentation.pptx"):
        os.remove("test_presentation.pptx")

    def test_presentation_creation(self):
        # Test if a presentation object is created successfully
        self.assertIsInstance(self.prs, pptx.Presentation)

    def test_title_slide_content(self):
        # Test if the title slide has the correct title and subtitle
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title_text = "Test Title"
        subtitle_text = "Test Subtitle"
        title.text = title_text
        subtitle.text = subtitle_text

        self.assertEqual(title.text, title_text)
        self.assertEqual(subtitle.text, subtitle_text)
        
        self.prs.save("test_presentation.pptx")

        prs2 = pptx.Presentation("test_presentation.pptx")
        self.assertEqual(prs2.slides[0].shapes.title.text, title_text)
        self.assertEqual(prs2.slides[0].placeholders[1].text, subtitle_text)
        

    def test_slide_layouts(self):
        # Test if the expected slide layouts exist
        self.assertGreater(len(self.prs.slide_layouts), 0)
        self.assertIn(0, range(len(self.prs.slide_layouts)))  # Test for Title Slide
        self.assertIn(1, range(len(self.prs.slide_layouts)))  # Test for Title and Content slide

    def test_add_bullet_point_slide(self):
        # Test adding a slide with bullet points
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title_text = "Test Bullet Slide"
        title.text = title_text

        placeholder = slide.placeholders[1]

        bullet_points = ["Bullet 1", "Bullet 2", "Bullet 3"]
        placeholder.text = bullet_points[0]
        for point in bullet_points[1:]:
          tf = placeholder.text_frame
          tf.add_paragraph().text=point

        self.assertEqual(title.text, title_text)
        self.assertEqual(placeholder.text_frame.paragraphs[0].text, bullet_points[0])
        self.assertEqual(placeholder.text_frame.paragraphs[1].text, bullet_points[1])
        self.assertEqual(placeholder.text_frame.paragraphs[2].text, bullet_points[2])

        self.prs.save("test_presentation.pptx")
        
        prs2 = pptx.Presentation("test_presentation.pptx")
        self.assertEqual(prs2.slides[0].shapes.title.text, title_text)
        self.assertEqual(prs2.slides[0].placeholders[1].text_frame.paragraphs[0].text, bullet_points[0])
        self.assertEqual(prs2.slides[0].placeholders[1].text_frame.paragraphs[1].text, bullet_points[1])
        self.assertEqual(prs2.slides[0].placeholders[1].text_frame.paragraphs[2].text, bullet_points[2])


    def test_save_presentation(self):
      # Test if the powerpoint file can be saved correctly
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test"
        subtitle.text = "Test"

        self.prs.save("test_presentation.pptx")
        self.assertTrue(os.path.exists("test_presentation.pptx"))
        self.assertGreater(os.path.getsize("test_presentation.pptx"), 0)

if __name__ == '__main__':
    unittest.main(argv=['first-arg-is-ignored'], exit=False)
